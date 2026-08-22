#!/usr/bin/env python3
"""Build the final project presentation from the existing teleoperation deck."""

from __future__ import annotations

import csv
from collections import Counter
from pathlib import Path

import h5py
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "IK_Teleoperation_v2.pptx"
OUTPUT = ROOT / "IK_Teleoperation_Final.pptx"
ASSET_DIR = ROOT / "outputs" / "presentation_assets"
SUMMARY_CSV = (
    ROOT
    / "outputs"
    / "evaluation"
    / "can_color_sort_pte_m005"
    / "experiment_summary.csv"
)
DATASET_DIR = ROOT / "datasets" / "can_color_sort"

FONT = "Noto Sans CJK KR"
FONT_EN = "Play"
NAVY = RGBColor(0x0F, 0x1B, 0x2D)
BLUE = RGBColor(0x1F, 0x5A, 0xA6)
BLUE_2 = RGBColor(0x3B, 0x82, 0xF6)
PALE_BLUE = RGBColor(0xE9, 0xF2, 0xFC)
PALE_NAVY = RGBColor(0xF3, 0xF6, 0xFA)
GRAY = RGBColor(0x58, 0x67, 0x7B)
MUTED = RGBColor(0x7A, 0x8A, 0xA0)
LIGHT = RGBColor(0xD8, 0xE1, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xD9, 0x3A, 0x32)
ORANGE = RGBColor(0xEF, 0x8A, 0x17)
GREEN = RGBColor(0x2C, 0xA0, 0x58)
CYAN = RGBColor(0x20, 0xA4, 0xB8)


def rgb_hex(color: RGBColor) -> str:
    return f"#{color}"


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: RGBColor = NAVY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font: str = FONT,
    margin: float = 0.03,
    line_spacing: float | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    paragraphs = text.split("\n")
    for idx, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return shape


def add_rich_text(
    slide, runs, x, y, w, h, *, size=16, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for value, color, bold in runs:
        run = p.add_run()
        run.text = value
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return shape


def add_box(
    slide,
    x,
    y,
    w,
    h,
    *,
    fill=WHITE,
    line=LIGHT,
    radius=True,
    line_width=1.0,
):
    kind = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def add_circle(slide, x, y, d, fill, line=WHITE, line_width=1.5):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=BLUE, width=2.2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_header(slide, title: str, number: int, section: str):
    add_text(slide, title, 0.55, 0.25, 10.9, 0.48, size=25.5, bold=True, font=FONT_EN)
    add_text(
        slide,
        section.upper(),
        10.5,
        0.34,
        2.25,
        0.24,
        size=8.5,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
        font=FONT_EN,
    )
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.57),
        Inches(0.84),
        Inches(12.18),
        Inches(0.025),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT
    line.line.fill.background()
    add_text(
        slide,
        f"{number:02d}",
        12.05,
        7.04,
        0.65,
        0.18,
        size=7.5,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        font=FONT_EN,
    )


def add_card(
    slide,
    title,
    body,
    x,
    y,
    w,
    h,
    *,
    accent=BLUE,
    fill=WHITE,
    title_size=15,
    body_size=11.5,
):
    add_box(slide, x, y, w, h, fill=fill, line=LIGHT)
    accent_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x + 0.16),
        Inches(y + 0.18),
        Inches(0.08),
        Inches(h - 0.36),
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.fill.background()
    add_text(
        slide,
        title,
        x + 0.36,
        y + 0.17,
        w - 0.5,
        0.32,
        size=title_size,
        bold=True,
        color=NAVY,
    )
    add_text(
        slide,
        body,
        x + 0.36,
        y + 0.62,
        w - 0.5,
        h - 0.75,
        size=body_size,
        color=GRAY,
        line_spacing=1.1,
    )


def add_media_placeholder(slide, label, hint, x, y, w, h, *, accent=BLUE):
    shape = add_box(
        slide, x, y, w, h, fill=PALE_NAVY, line=accent, radius=False, line_width=1.4
    )
    shape.line.dash_style = 4
    add_circle(slide, x + w / 2 - 0.26, y + h / 2 - 0.62, 0.52, accent)
    add_text(
        slide,
        "+",
        x + w / 2 - 0.26,
        y + h / 2 - 0.54,
        0.52,
        0.22,
        size=14,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        font=FONT_EN,
    )
    add_text(
        slide,
        label,
        x + 0.20,
        y + h / 2 + 0.05,
        w - 0.40,
        0.30,
        size=13,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
        font=FONT_EN,
    )
    add_text(
        slide,
        hint,
        x + 0.24,
        y + h / 2 + 0.43,
        w - 0.48,
        0.48,
        size=10.5,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    return shape


def add_picture_fit(slide, path: Path, x, y, w, h, *, line=LIGHT):
    with Image.open(path) as im:
        iw, ih = im.size
    target = w / h
    source = iw / ih
    if source > target:
        pic_h = h
        pic_w = h * source
        px = x - (pic_w - w) / 2
        py = y
    else:
        pic_w = w
        pic_h = w / source
        px = x
        py = y - (pic_h - h) / 2
    frame = add_box(slide, x, y, w, h, fill=WHITE, line=line, radius=False)
    pic = slide.shapes.add_picture(
        str(path), Inches(px), Inches(py), width=Inches(pic_w), height=Inches(pic_h)
    )
    # Crop by relying on a white cover rectangle only when needed is awkward in pptx;
    # use native crop ratios so the image stays inside the requested frame.
    if source > target:
        visible_ratio = target / source
        crop = (1 - visible_ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
        pic.left, pic.top, pic.width, pic.height = (
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
    elif source < target:
        visible_ratio = source / target
        crop = (1 - visible_ratio) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
        pic.left, pic.top, pic.width, pic.height = (
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
    frame.element.addnext(pic.element)
    return pic


def read_results():
    with SUMMARY_CSV.open(newline="", encoding="utf-8-sig") as f:
        rows = list(csv.DictReader(f))
    for row in rows:
        for key in (
            "data_count",
            "pte_steps",
            "trials",
            "success_rate_pct",
            "success_ci95_low_pct",
            "success_ci95_high_pct",
            "median_success_time_s",
            "mean_penalized_time_s",
            "speedup_vs_f0",
        ):
            row[key] = float(row[key]) if row[key] else np.nan
    return rows


def scan_dataset():
    counts = Counter()
    samples = {}
    first = None
    for path in sorted(DATASET_DIR.glob("episode_*.hdf5")):
        with h5py.File(path, "r") as h5:
            variant = str(h5.attrs.get("object_variant", "unknown"))
            counts[variant] += 1
            samples.setdefault(variant, path)
            if first is None:
                first = path
    return counts, samples, first


def build_assets(rows, counts, samples, collector_episode):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "axes.edgecolor": "#D8E1EC",
            "axes.labelcolor": "#58677B",
            "xtick.color": "#58677B",
            "ytick.color": "#58677B",
            "axes.titlecolor": "#0F1B2D",
        }
    )

    palette = {
        (97, "joint"): "#7A8AA0",
        (97, "task"): "#20A4B8",
        (150, "joint"): "#1F5AA6",
        (150, "task"): "#EF8A17",
    }
    labels = {
        (97, "joint"): "D97 · Joint",
        (97, "task"): "D97 · Task",
        (150, "joint"): "D150 · Joint",
        (150, "task"): "D150 · Task",
    }

    fig, ax = plt.subplots(figsize=(9.2, 4.6), dpi=180)
    for key in labels:
        subset = sorted(
            (
                r
                for r in rows
                if int(r["data_count"]) == key[0] and r["representation"] == key[1]
            ),
            key=lambda r: r["pte_steps"],
        )
        x = np.array([r["pte_steps"] for r in subset])
        y = np.array([r["success_rate_pct"] for r in subset])
        low = np.array([r["success_ci95_low_pct"] for r in subset])
        high = np.array([r["success_ci95_high_pct"] for r in subset])
        ax.errorbar(
            x,
            y,
            yerr=[y - low, high - y],
            marker="o",
            linewidth=2.3,
            markersize=5.5,
            capsize=3,
            label=labels[key],
            color=palette[key],
        )
    ax.axvspan(-0.5, 10.5, color="#E9F2FC", alpha=0.65, zorder=-2)
    ax.axvline(10.5, color="#D93A32", linestyle="--", linewidth=1.2, alpha=0.65)
    ax.text(10.9, 12, "reliability cliff", color="#D93A32", fontsize=9)
    ax.set(
        xlabel="PTE future step f",
        ylabel="Success rate (%)",
        xticks=[0, 5, 10, 15, 20],
        ylim=(-3, 105),
    )
    ax.grid(axis="y", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, ncol=2, loc="lower left", fontsize=8.5)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "success_rate.png", transparent=False, facecolor="white")
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(9.2, 4.4), dpi=180)
    for key in labels:
        subset = sorted(
            (
                r
                for r in rows
                if int(r["data_count"]) == key[0] and r["representation"] == key[1]
            ),
            key=lambda r: r["pte_steps"],
        )
        ax.plot(
            [r["pte_steps"] for r in subset],
            [r["mean_penalized_time_s"] for r in subset],
            marker="o",
            linewidth=2.3,
            markersize=5.5,
            label=labels[key],
            color=palette[key],
        )
    ax.axhline(20, color="#D8E1EC", linewidth=1)
    ax.axvspan(-0.5, 10.5, color="#E9F2FC", alpha=0.65, zorder=-2)
    ax.set(
        xlabel="PTE future step f",
        ylabel="Penalized completion time (s)",
        xticks=[0, 5, 10, 15, 20],
        ylim=(6.5, 20.8),
    )
    ax.grid(axis="y", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, ncol=2, loc="upper left", fontsize=8.5)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "penalized_time.png", facecolor="white")
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(5.1, 4.4), dpi=180)
    for key in labels:
        subset = [
            r
            for r in rows
            if int(r["data_count"]) == key[0] and r["representation"] == key[1]
        ]
        ax.scatter(
            [r["mean_penalized_time_s"] for r in subset],
            [r["success_rate_pct"] for r in subset],
            s=48,
            color=palette[key],
            label=labels[key],
            alpha=0.92,
        )
        for r in subset:
            ax.annotate(
                f"f{int(r['pte_steps'])}",
                (r["mean_penalized_time_s"], r["success_rate_pct"]),
                xytext=(4, 3),
                textcoords="offset points",
                fontsize=7,
                color=palette[key],
            )
    ax.set(
        xlabel="Penalized time (s)  ← faster",
        ylabel="Success rate (%)",
        xlim=(7, 20.5),
        ylim=(-3, 105),
    )
    ax.grid(alpha=0.2)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "pareto.png", facecolor="white")
    plt.close(fig)

    variants = ["green", "red", "orange", "blue"]
    base = [50, 47, 0, 0]
    aug = [counts.get(v, 0) for v in variants]
    colors = ["#2CA058", "#D93A32", "#EF8A17", "#1F5AA6"]
    fig, ax = plt.subplots(figsize=(6.6, 3.4), dpi=180)
    bottoms = np.zeros(2)
    for i, v in enumerate(variants):
        vals = [base[i], aug[i]]
        ax.bar(
            [0, 1],
            vals,
            bottom=bottoms,
            width=0.56,
            label=v.capitalize(),
            color=colors[i],
        )
        for xi, val, bottom in zip([0, 1], vals, bottoms):
            if val:
                ax.text(
                    xi,
                    bottom + val / 2,
                    str(val),
                    ha="center",
                    va="center",
                    color="white",
                    fontsize=10,
                    fontweight="bold",
                )
        bottoms += vals
    ax.set_xticks([0, 1], ["D97\nbase", "D150\naugmented"])
    ax.set_ylabel("Episodes")
    ax.set_ylim(0, 160)
    ax.grid(axis="y", alpha=0.18)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(
        frameon=False,
        ncol=4,
        loc="upper center",
        bbox_to_anchor=(0.5, 1.13),
        fontsize=8,
    )
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "dataset_distribution.png", facecolor="white")
    plt.close(fig)


def add_cover(prs):
    blank = next(
        (layout for layout in prs.slide_layouts if layout.name.lower() == "blank"),
        prs.slide_layouts[-1],
    )
    slide = prs.slides.add_slide(blank)
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    add_text(
        slide,
        "FINAL PROJECT",
        0.72,
        0.72,
        3.0,
        0.26,
        size=10,
        color=BLUE,
        bold=True,
        font=FONT_EN,
    )
    add_text(
        slide,
        "FFW-SH5 기반 모방학습\n파이프라인 구축 및 분석",
        0.72,
        1.28,
        6.15,
        1.38,
        size=29,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        "MuJoCo Teleoperation · ACT · Color Sorting",
        0.74,
        2.92,
        5.95,
        0.36,
        size=15.5,
        color=GRAY,
        bold=True,
        font=FONT_EN,
    )
    add_text(
        slide,
        "JOINT SPACE  ↔  TASK SPACE   |   PTE",
        0.74,
        3.48,
        5.60,
        0.30,
        size=11.5,
        color=BLUE,
        bold=True,
        font=FONT_EN,
    )
    add_box(slide, 0.74, 4.45, 5.58, 0.78, fill=PALE_BLUE, line=PALE_BLUE)
    add_text(
        slide,
        "데이터 취득부터 2,000회 정책 평가까지",
        1.02,
        4.70,
        5.03,
        0.26,
        size=13.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "2026.08.23  윤우성",
        0.75,
        6.62,
        2.65,
        0.24,
        size=10.5,
        color=MUTED,
        bold=True,
    )
    add_media_placeholder(
        slide,
        "TITLE IMAGE / VIDEO",
        "대표 시뮬레이션 장면 또는 데모 썸네일",
        7.18,
        0.72,
        5.48,
        6.08,
        accent=BLUE,
    )
    return slide


def append_main_slides(prs, rows, counts):
    blank = next(
        (layout for layout in prs.slide_layouts if layout.name.lower() == "blank"),
        prs.slide_layouts[-1],
    )
    new_slides = []

    # 02 — problem definition
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "문제 정의와 연구 질문", 2, "Research Question")
    add_text(
        s,
        "색을 인식하고, 오른팔로 캔을 집어\n올바른 상자에 분류할 수 있는가?",
        0.65,
        1.12,
        6.3,
        0.95,
        size=23,
        bold=True,
    )
    add_rich_text(
        s,
        [
            ("Warm ", RED, True),
            ("red · orange", RED, True),
            ("  →  RED BIN", NAVY, True),
        ],
        7.35,
        1.28,
        5.1,
        0.36,
        size=15,
    )
    add_rich_text(
        s,
        [
            ("Cool ", BLUE, True),
            ("green · blue", BLUE, True),
            ("  →  BLUE BIN", NAVY, True),
        ],
        7.35,
        1.77,
        5.1,
        0.36,
        size=15,
    )
    for i, (c, label) in enumerate(
        [(RED, "R"), (ORANGE, "O"), (GREEN, "G"), (BLUE, "B")]
    ):
        add_circle(s, 8.05 + i * 0.86, 2.40, 0.48, c)
        add_text(
            s,
            label,
            8.05 + i * 0.86,
            2.48,
            0.48,
            0.18,
            size=9,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    add_arrow(s, 8.76, 3.04, 8.33, 3.72, RED)
    add_arrow(s, 10.38, 3.04, 11.02, 3.72, BLUE)
    add_box(s, 7.62, 3.72, 1.65, 1.18, fill=RED, line=RED)
    add_box(s, 10.27, 3.72, 1.65, 1.18, fill=BLUE, line=BLUE)
    add_text(
        s,
        "RED",
        7.62,
        4.10,
        1.65,
        0.25,
        size=13,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        "BLUE",
        10.27,
        4.10,
        1.65,
        0.25,
        size=13,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_card(
        s,
        "Representation",
        "Joint space와 Task space 중\n어떤 출력 표현이 유리한가?",
        0.65,
        2.45,
        2.06,
        2.04,
        accent=BLUE,
    )
    add_card(
        s,
        "Data",
        "데이터 수와 색상 다양성이\n성공률에 어떤 영향을 주는가?",
        2.89,
        2.45,
        2.06,
        2.04,
        accent=GREEN,
    )
    add_card(
        s,
        "Inference",
        "미래 action을 앞당겨 쓰는 PTE가\n속도와 안정성을 개선하는가?",
        5.13,
        2.45,
        2.06,
        2.04,
        accent=ORANGE,
    )
    add_text(
        s,
        "IL로 end-to-end 색상 분류 정책을 학습하고 세 축을 통제 실험",
        0.67,
        5.20,
        11.8,
        0.42,
        size=18,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )

    # 03 — constraints
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "제약 조건을 시스템 설계로 전환", 3, "Approach")
    items = [
        (
            "로봇 하드웨어 부재",
            "MuJoCo 기반 FFW-SH5 디지털 트윈",
            "충돌·물성·카메라를 포함한 반복 가능한 환경",
            BLUE,
        ),
        (
            "직접 조작 장치 부재",
            "6-DoF 목표 pose + IK teleoperation",
            "마우스/키보드 입력을 관절 명령으로 변환",
            CYAN,
        ),
        (
            "비교 가능한 학습 필요",
            "모듈형 데이터·정책·평가 파이프라인",
            "env / representation / PTE를 실행 인자로 분리",
            ORANGE,
        ),
    ]
    for i, (problem, solution, detail, accent) in enumerate(items):
        y = 1.2 + i * 1.63
        add_box(s, 0.7, y, 3.05, 1.1, fill=PALE_NAVY, line=LIGHT)
        add_text(
            s,
            "CONSTRAINT",
            0.95,
            y + 0.16,
            1.15,
            0.18,
            size=8,
            color=MUTED,
            bold=True,
            font=FONT_EN,
        )
        add_text(s, problem, 0.95, y + 0.45, 2.55, 0.31, size=16, bold=True)
        add_arrow(s, 3.92, y + 0.55, 4.72, y + 0.55, accent, 2.8)
        add_box(s, 4.9, y - 0.08, 7.65, 1.28, fill=WHITE, line=accent, line_width=1.4)
        add_text(
            s, solution, 5.2, y + 0.12, 6.95, 0.34, size=17, bold=True, color=accent
        )
        add_text(s, detail, 5.2, y + 0.58, 6.95, 0.31, size=12.5, color=GRAY)
    add_text(
        s,
        "결과: 데이터 취득 → 학습 → 시뮬레이션 평가까지 하나의 재현 가능한 파이프라인",
        0.7,
        6.28,
        11.85,
        0.36,
        size=16,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )

    # 04 — pipeline
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "전체 시스템 파이프라인", 4, "System")
    stages = [
        ("01", "Scenario", "모듈형 환경\n색·상자 랜덤화"),
        ("02", "Teleoperation", "6-DoF 목표\nIK + grasp"),
        ("03", "Dataset", "HDF5 · 25 Hz\nimage/state/action"),
        ("04", "ACT Training", "Joint / Task\n동일 hyperparameter"),
        ("05", "Inference", "Temporal ensemble\nPTE f-step"),
        ("06", "Evaluation", "20 조건 × 100회\n성공률·시간"),
    ]
    for i, (num, title, body) in enumerate(stages):
        x = 0.52 + i * 2.10
        fill = PALE_BLUE if i in (3, 4) else WHITE
        add_box(s, x, 2.02, 1.77, 2.17, fill=fill, line=BLUE if i in (3, 4) else LIGHT)
        add_circle(s, x + 0.56, 1.45, 0.64, BLUE if i < 4 else ORANGE)
        add_text(
            s,
            num,
            x + 0.56,
            1.61,
            0.64,
            0.18,
            size=9,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            font=FONT_EN,
        )
        add_text(
            s,
            title,
            x + 0.13,
            2.33,
            1.51,
            0.32,
            size=14.5,
            bold=True,
            color=NAVY,
            align=PP_ALIGN.CENTER,
            font=FONT_EN,
        )
        add_text(
            s,
            body,
            x + 0.14,
            2.93,
            1.49,
            0.68,
            size=11.2,
            color=GRAY,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if i < len(stages) - 1:
            add_arrow(s, x + 1.79, 3.05, x + 2.04, 3.05, MUTED, 1.5)
    add_box(s, 1.2, 5.10, 10.9, 0.95, fill=PALE_NAVY, line=LIGHT)
    add_text(
        s, "핵심 설계 원칙", 1.48, 5.35, 1.55, 0.25, size=13.5, bold=True, color=BLUE
    )
    add_text(
        s,
        "환경은 --env, 정책 표현은 --policy-representation, PTE는 --policy-pte-steps로 분리",
        3.12,
        5.31,
        8.53,
        0.34,
        size=14,
        color=NAVY,
    )

    # 05 — teleop and IK
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "MuJoCo Teleoperation과 IK", 5, "Data Acquisition")
    add_media_placeholder(
        s,
        "VIDEO / IMAGE",
        "MuJoCo teleoperation 전체 장면 또는 10–15초 데모",
        0.62,
        1.14,
        6.35,
        4.77,
        accent=BLUE,
    )
    flow = [
        ("Operator", "6-DoF target pose\n+ grasp command"),
        ("IK Solver", "EE error → Jacobian\ncollision-aware update"),
        ("MuJoCo", "q target → physics\n25 Hz closed loop"),
    ]
    for i, (title, body) in enumerate(flow):
        y = 1.22 + i * 1.42
        add_box(
            s,
            7.45,
            y,
            4.95,
            1.0,
            fill=PALE_BLUE if i == 1 else WHITE,
            line=BLUE if i == 1 else LIGHT,
        )
        add_text(
            s,
            title,
            7.72,
            y + 0.15,
            1.28,
            0.26,
            size=14,
            bold=True,
            color=BLUE,
            font=FONT_EN,
        )
        add_text(
            s,
            body,
            9.00,
            y + 0.13,
            3.08,
            0.55,
            size=11.5,
            color=GRAY,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if i < 2:
            add_arrow(s, 9.91, y + 1.02, 9.91, y + 1.35, BLUE, 1.8)
    add_card(
        s,
        "Why IK?",
        "Task-space 시연을 joint trajectory로 실행하면서\n동시에 EE pose를 데이터에 저장",
        7.45,
        5.45,
        4.95,
        0.88,
        accent=ORANGE,
        body_size=10.8,
        title_size=13,
    )

    # 06 — collector
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "Data Collector: 하나의 시연, 두 표현 공간", 6, "Dataset")
    add_media_placeholder(
        s, "cam_high", "외부 카메라 프레임", 5.65, 1.22, 2.06, 2.18, accent=BLUE
    )
    add_media_placeholder(
        s, "cam_left_wrist", "왼손목 카메라 프레임", 7.99, 1.22, 2.06, 2.18, accent=CYAN
    )
    add_media_placeholder(
        s,
        "cam_right_wrist",
        "오른손목 카메라 프레임",
        10.33,
        1.22,
        2.06,
        2.18,
        accent=ORANGE,
    )
    add_box(s, 0.68, 1.22, 4.53, 4.90, fill=RGBColor(0x13, 0x22, 0x38), line=NAVY)
    schema = (
        "episode_xxxxxx.hdf5\n"
        "├─ observations/images\n"
        "│  ├─ cam_high\n"
        "│  ├─ cam_left_wrist\n"
        "│  └─ cam_right_wrist\n"
        "├─ observations/qpos, qvel\n"
        "├─ observations/ee_pose\n"
        "│  ├─ left  [x y z qw qx qy qz]\n"
        "│  └─ right [x y z qw qx qy qz]\n"
        "├─ action\n"
        "└─ attrs: success, color, target, seed"
    )
    add_text(
        s,
        schema,
        0.96,
        1.52,
        4.05,
        4.35,
        size=11.3,
        color=WHITE,
        font="DejaVu Sans Mono",
        line_spacing=1.05,
    )
    add_card(
        s,
        "Joint policy",
        "right q(7) + grasp(1)",
        5.66,
        4.24,
        3.17,
        1.22,
        accent=BLUE,
        fill=PALE_BLUE,
    )
    add_card(
        s,
        "Task policy",
        "right EE pose(7) + grasp(1)",
        9.11,
        4.24,
        3.30,
        1.22,
        accent=ORANGE,
        fill=RGBColor(0xFF, 0xF6, 0xE8),
    )
    add_text(
        s,
        "수집 단계에서 모두 저장 → 학습 설정에서 필요한 representation만 선택",
        5.70,
        5.83,
        6.65,
        0.31,
        size=13,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )

    # 07 — dataset evolution
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "Dataset Evolution: 수량 + 색상 다양성", 7, "Dataset")
    add_picture_fit(s, ASSET_DIR / "dataset_distribution.png", 0.65, 1.15, 5.38, 3.12)
    montage_colors = [
        ("GREEN", GREEN),
        ("RED", RED),
        ("ORANGE", ORANGE),
        ("BLUE", BLUE),
    ]
    for i, (label, accent) in enumerate(montage_colors):
        add_media_placeholder(
            s, label, "대표 episode", 6.35 + i * 1.51, 1.18, 1.36, 2.05, accent=accent
        )
    add_card(
        s,
        "D97 · Base",
        "Green 50 + Red 47\n초기 약 100개 중 실패 시연 제거",
        0.78,
        4.52,
        4.8,
        1.27,
        accent=GRAY,
    )
    add_card(
        s,
        "D150 · Augmented",
        "D97 + Orange 24 + Blue 29\n총 150 episodes, 색상 분포 확장",
        6.35,
        4.52,
        5.95,
        1.27,
        accent=BLUE,
        fill=PALE_BLUE,
    )
    add_text(
        s,
        "주의: D97↔D150 비교는 데이터 개수뿐 아니라 색상 다양성도 함께 변한 조건",
        0.72,
        6.25,
        11.6,
        0.34,
        size=13.5,
        bold=True,
        color=RED,
        align=PP_ALIGN.CENTER,
    )

    # 08 — ACT
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "ACT 학습 파이프라인", 8, "Model")
    components = [
        ("Images", "cam_high\ncam_right_wrist", BLUE),
        ("State", "q(7)+grasp\nor pose(7)+grasp", CYAN),
        ("Encoder", "ResNet18\nvisual features", GREEN),
        ("ACT", "CVAE Transformer\naction chunk = 90", ORANGE),
        ("Controller", "ensemble / PTE\n25 Hz execution", RED),
    ]
    for i, (title, body, accent) in enumerate(components):
        x = 0.53 + i * 2.53
        add_box(
            s,
            x,
            2.02,
            2.10,
            2.10,
            fill=WHITE if i != 3 else RGBColor(0xFF, 0xF6, 0xE8),
            line=accent,
            line_width=1.4,
        )
        add_text(
            s,
            title,
            x + 0.18,
            2.38,
            1.74,
            0.29,
            size=16,
            bold=True,
            color=accent,
            align=PP_ALIGN.CENTER,
            font=FONT_EN,
        )
        add_text(
            s,
            body,
            x + 0.14,
            3.02,
            1.82,
            0.62,
            size=11,
            color=GRAY,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if i < 4:
            add_arrow(s, x + 2.12, 3.06, x + 2.45, 3.06, MUTED, 1.6)
    add_box(s, 1.05, 5.03, 11.15, 0.86, fill=PALE_NAVY, line=LIGHT)
    add_rich_text(
        s,
        [
            ("통제 변수  ", MUTED, True),
            ("동일한 ACT 구조와 주요 hyperparameter", NAVY, True),
            ("  ·  변경 변수  ", MUTED, True),
            ("representation / dataset", BLUE, True),
        ],
        1.34,
        5.26,
        10.6,
        0.31,
        size=14,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        "Chunk 단위 예측으로 장기 동작을 생성하고, 실행 시 겹치는 예측을 결합",
        1.08,
        6.20,
        11.1,
        0.32,
        size=13,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    # 09 — representation comparison
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "Joint Space vs Task Space", 9, "Representation")
    add_box(s, 0.66, 1.20, 5.85, 4.95, fill=PALE_BLUE, line=BLUE, line_width=1.4)
    add_box(
        s,
        6.82,
        1.20,
        5.85,
        4.95,
        fill=RGBColor(0xFF, 0xF6, 0xE8),
        line=ORANGE,
        line_width=1.4,
    )
    add_text(
        s,
        "JOINT SPACE",
        1.03,
        1.54,
        5.10,
        0.39,
        size=22,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
        font=FONT_EN,
    )
    add_text(
        s,
        "TASK SPACE",
        7.18,
        1.54,
        5.10,
        0.39,
        size=22,
        bold=True,
        color=ORANGE,
        align=PP_ALIGN.CENTER,
        font=FONT_EN,
    )
    add_text(
        s,
        "state / action",
        1.12,
        2.24,
        1.30,
        0.26,
        size=10,
        color=MUTED,
        bold=True,
        font=FONT_EN,
    )
    add_text(
        s,
        "[ q₁ … q₇, grasp ]  ∈ ℝ⁸",
        2.38,
        2.17,
        3.43,
        0.38,
        size=17,
        bold=True,
        color=NAVY,
        font="DejaVu Sans",
    )
    add_text(
        s,
        "state / action",
        7.28,
        2.24,
        1.30,
        0.26,
        size=10,
        color=MUTED,
        bold=True,
        font=FONT_EN,
    )
    add_text(
        s,
        "[ x y z, qʷ qˣ qʸ qᶻ, grasp ] ∈ ℝ⁸",
        8.52,
        2.15,
        3.73,
        0.47,
        size=15,
        bold=True,
        color=NAVY,
        font="DejaVu Sans",
    )
    joint_rows = [
        ("Execution", "policy output → joint target"),
        ("Strength", "IK 오차 없이 직접 실행"),
        ("Risk", "기구학 구조에 종속"),
    ]
    task_rows = [
        ("Execution", "policy output → IK → joint target"),
        ("Strength", "작업 의미와 직접 대응"),
        ("Risk", "IK 추종 오차·충돌 제약"),
    ]
    for i, ((a, b), (c, d)) in enumerate(zip(joint_rows, task_rows)):
        y = 3.08 + i * 0.77
        add_text(
            s, a, 1.08, y, 1.08, 0.24, size=10.5, color=BLUE, bold=True, font=FONT_EN
        )
        add_text(s, b, 2.30, y - 0.02, 3.78, 0.30, size=12.2, color=GRAY)
        add_text(
            s, c, 7.25, y, 1.08, 0.24, size=10.5, color=ORANGE, bold=True, font=FONT_EN
        )
        add_text(s, d, 8.48, y - 0.02, 3.77, 0.30, size=12.2, color=GRAY)
    add_text(
        s,
        "같은 이미지 입력과 ACT 구조에서 표현 공간만 교체",
        0.72,
        6.42,
        11.85,
        0.31,
        size=14,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )

    # 10 — PTE
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "PTE: 미래 action을 선택하는 추론 제어", 10, "Inference")
    add_text(
        s,
        "ACT Temporal Ensemble",
        0.75,
        1.15,
        5.45,
        0.34,
        size=17,
        bold=True,
        color=BLUE,
        font=FONT_EN,
    )
    add_text(
        s,
        "현재 시점에 해당하는 예측을 여러 chunk에서 가중 평균",
        0.76,
        1.57,
        5.52,
        0.32,
        size=12.5,
        color=GRAY,
    )
    add_text(
        s,
        "Proleptic Temporal Ensemble",
        6.78,
        1.15,
        5.45,
        0.34,
        size=17,
        bold=True,
        color=ORANGE,
        font=FONT_EN,
    )
    add_text(
        s,
        "f-step 미래 action들을 결합해 동작을 선행",
        6.79,
        1.57,
        5.52,
        0.32,
        size=12.5,
        color=GRAY,
    )
    cell_colors = [
        RGBColor(0xF5, 0xD6, 0xC9),
        RGBColor(0xE7, 0xCC, 0xE8),
        RGBColor(0xFF, 0xEF, 0x59),
        RGBColor(0x89, 0xB2, 0xE8),
        RGBColor(0xB8, 0xE6, 0xBE),
    ]
    for block, base_x, offset in [(0, 0.96, 0), (1, 7.00, 2)]:
        for r in range(5):
            for c in range(5):
                x = base_x + (c + r * 0.46) * 0.64
                y = 2.18 + r * 0.60
                add_box(
                    s,
                    x,
                    y,
                    0.57,
                    0.47,
                    fill=cell_colors[c],
                    line=NAVY,
                    radius=True,
                    line_width=0.8,
                )
        hx = base_x + (offset + 2 * 0.46) * 0.64 - 0.10
        h = add_box(
            s,
            hx,
            2.08,
            0.78,
            3.05,
            fill=WHITE,
            line=ORANGE if block else RED,
            radius=True,
            line_width=2.0,
        )
        h.fill.transparency = 100
    add_text(
        s,
        "t−4\nt−3\nt−2\nt−1\n현재",
        0.25,
        2.20,
        0.58,
        2.90,
        size=11,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
        line_spacing=1.25,
    )
    add_text(
        s,
        "t−4\nt−3\nt−2\nt−1\n현재",
        6.28,
        2.20,
        0.58,
        2.90,
        size=11,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
        line_spacing=1.25,
    )
    add_box(s, 2.10, 5.60, 9.13, 0.64, fill=PALE_NAVY, line=LIGHT)
    add_rich_text(
        s,
        [
            ("f = 0, 5, 10, 15, 20", BLUE, True),
            ("  |  25 Hz 기준 look-ahead = ", GRAY, False),
            ("0.0–0.8 s", ORANGE, True),
            ("  |  재학습 불필요", NAVY, True),
        ],
        2.30,
        5.77,
        8.73,
        0.27,
        size=13,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        "가설: 적당한 f는 지연을 줄이지만, 과도한 f는 접촉 구간을 건너뛴다",
        1.10,
        6.45,
        11.1,
        0.30,
        size=14,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )

    # 11 — experiment design
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "실험 설계", 11, "Evaluation")
    add_text(
        s,
        "3개 실험축을 완전 요인 설계로 조합",
        0.68,
        1.12,
        6.9,
        0.40,
        size=20,
        bold=True,
    )
    factors = [
        ("Representation", "2", "Joint / Task", BLUE),
        ("Dataset", "2", "D97 / D150", GREEN),
        ("PTE future step", "5", "0 / 5 / 10 / 15 / 20", ORANGE),
        ("Rollout", "100", "per condition", RED),
    ]
    for i, (title, n, body, accent) in enumerate(factors):
        x = 0.70 + i * 3.04
        add_box(s, x, 1.90, 2.68, 1.62, fill=WHITE, line=accent, line_width=1.3)
        add_text(
            s,
            title,
            x + 0.16,
            2.12,
            2.36,
            0.24,
            size=10.5,
            color=accent,
            bold=True,
            align=PP_ALIGN.CENTER,
            font=FONT_EN,
        )
        add_text(
            s,
            n,
            x + 0.20,
            2.51,
            2.28,
            0.47,
            size=27,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
            font=FONT_EN,
        )
        add_text(
            s,
            body,
            x + 0.16,
            3.06,
            2.36,
            0.20,
            size=10.5,
            color=GRAY,
            align=PP_ALIGN.CENTER,
            font=FONT_EN,
        )
    add_text(
        s,
        "2 × 2 × 5 × 100  =  2,000 rollouts",
        1.38,
        4.08,
        10.58,
        0.58,
        size=26,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
        font=FONT_EN,
    )
    add_box(s, 0.75, 5.15, 11.8, 1.05, fill=PALE_NAVY, line=LIGHT)
    add_text(
        s,
        "Metrics",
        1.04,
        5.43,
        1.15,
        0.26,
        size=12.5,
        color=MUTED,
        bold=True,
        font=FONT_EN,
    )
    add_text(
        s,
        "성공률 (95% CI)  ·  성공시간 중앙값  ·  실패=20 s 패널티 시간  ·  f=0 대비 속도향상",
        2.20,
        5.38,
        9.85,
        0.34,
        size=13.2,
        color=NAVY,
    )
    add_text(
        s,
        "모든 조건은 같은 task와 평가 프로토콜 사용",
        0.76,
        6.52,
        11.75,
        0.28,
        size=12.5,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    # 12 — success results
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "결과 1: 성공률", 12, "Results")
    add_picture_fit(s, ASSET_DIR / "success_rate.png", 0.60, 1.08, 8.28, 4.95)
    add_card(
        s,
        "Best reliability",
        "D150 · Joint\nf = 0 / 5 → 100%",
        9.20,
        1.23,
        3.05,
        1.33,
        accent=BLUE,
        fill=PALE_BLUE,
        body_size=12.5,
    )
    add_card(
        s,
        "Task strength",
        "D97 · Task\nf = 5 → 100%",
        9.20,
        2.82,
        3.05,
        1.33,
        accent=CYAN,
        body_size=12.5,
    )
    add_card(
        s,
        "Cliff",
        "f ≥ 15에서\n특히 Task 급락",
        9.20,
        4.41,
        3.05,
        1.33,
        accent=RED,
        fill=RGBColor(0xFF, 0xF0, 0xEF),
        body_size=12.5,
    )
    add_text(
        s,
        "f=5는 네 정책 모두에서 가장 안전한 공통 운용점",
        0.84,
        6.35,
        11.55,
        0.35,
        size=16,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )

    # 13 — speed and Pareto
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "결과 2: 속도–신뢰성 Trade-off", 13, "Results")
    add_picture_fit(s, ASSET_DIR / "penalized_time.png", 0.55, 1.13, 7.20, 4.65)
    add_picture_fit(s, ASSET_DIR / "pareto.png", 7.95, 1.13, 4.72, 4.65)
    add_box(s, 0.76, 6.02, 11.65, 0.66, fill=PALE_NAVY, line=LIGHT)
    add_rich_text(
        s,
        [
            ("최대 속도향상  ", MUTED, True),
            ("D97 Task · f10 = 1.507×", CYAN, True),
            ("   |   ", MUTED, False),
            ("고신뢰 고속  ", MUTED, True),
            ("D150 Joint · f10 = 98%, 1.424×", BLUE, True),
        ],
        0.98,
        6.21,
        11.2,
        0.26,
        size=13,
        align=PP_ALIGN.CENTER,
    )

    # 14 — discussion
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "해석과 핵심 발견", 14, "Discussion")
    findings = [
        (
            "01",
            "Data diversity matters",
            "D150 Joint f0: 100%\nD97 Joint f0: 79%",
            BLUE,
        ),
        (
            "02",
            "f=5 is robust",
            "네 정책에서 성공률 유지 또는 향상\n공통 기본값으로 적합",
            GREEN,
        ),
        (
            "03",
            "f=10 is aggressive",
            "최대 1.507× 빨라지지만\n정책별 성공률 손실 가능",
            ORANGE,
        ),
        (
            "04",
            "Representation interacts",
            "Task는 D97에서 강했지만\nD150·큰 f에서는 민감",
            RED,
        ),
    ]
    for i, (num, title, body, accent) in enumerate(findings):
        x = 0.69 + (i % 2) * 6.14
        y = 1.20 + (i // 2) * 2.25
        add_box(s, x, y, 5.72, 1.88, fill=WHITE, line=accent, line_width=1.3)
        add_circle(s, x + 0.25, y + 0.28, 0.55, accent)
        add_text(
            s,
            num,
            x + 0.25,
            y + 0.42,
            0.55,
            0.16,
            size=8.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            font=FONT_EN,
        )
        add_text(
            s,
            title,
            x + 1.00,
            y + 0.26,
            4.34,
            0.30,
            size=16,
            bold=True,
            color=NAVY,
            font=FONT_EN,
        )
        add_text(s, body, x + 1.00, y + 0.83, 4.35, 0.59, size=12.3, color=GRAY)
    add_box(s, 1.05, 5.83, 11.2, 0.66, fill=PALE_BLUE, line=BLUE)
    add_text(
        s,
        "결론: PTE는 정책을 재학습하지 않고 속도를 높이지만, f는 표현 공간별로 검증해야 한다",
        1.30,
        6.03,
        10.68,
        0.25,
        size=14,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )

    # 15 — conclusion and future
    s = prs.slides.add_slide(blank)
    new_slides.append(s)
    add_header(s, "Conclusion & Next Step", 15, "Closing")
    add_text(
        s,
        "시뮬레이션 안에서 IL 전체 루프를 완성",
        0.72,
        1.12,
        7.6,
        0.46,
        size=23,
        bold=True,
    )
    done = [
        "MuJoCo + IK 기반 teleoperation",
        "Joint/Task를 함께 담는 HDF5 collector",
        "ACT modular training & inference",
        "2,000회 정량 평가와 PTE 분석",
    ]
    for i, item in enumerate(done):
        y = 1.92 + i * 0.68
        add_circle(s, 0.86, y + 0.03, 0.28, BLUE)
        add_text(
            s,
            "✓",
            0.86,
            y + 0.065,
            0.28,
            0.14,
            size=8,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(s, item, 1.33, y - 0.02, 5.62, 0.30, size=14, color=NAVY)
    add_box(s, 7.50, 1.16, 4.83, 4.42, fill=PALE_NAVY, line=LIGHT)
    add_text(
        s,
        "NEXT",
        7.86,
        1.52,
        1.13,
        0.24,
        size=10,
        color=ORANGE,
        bold=True,
        font=FONT_EN,
    )
    nexts = [
        ("01", "Real robot deployment", "OpenArm / 보유 로봇에 sim-to-real"),
        ("02", "Robustness", "카메라·동역학·마찰 domain randomization"),
        ("03", "Adaptive PTE", "접촉 단계에 따라 f를 동적으로 선택"),
        ("04", "Policy comparison", "Diffusion Policy와 동일 protocol 비교"),
    ]
    for i, (num, title, body) in enumerate(nexts):
        y = 1.95 + i * 0.84
        add_text(
            s, num, 7.86, y, 0.46, 0.22, size=9, color=ORANGE, bold=True, font=FONT_EN
        )
        add_text(
            s,
            title,
            8.38,
            y - 0.03,
            3.42,
            0.25,
            size=12.5,
            bold=True,
            color=NAVY,
            font=FONT_EN,
        )
        add_text(s, body, 8.38, y + 0.28, 3.52, 0.27, size=10.2, color=GRAY)
    add_text(
        s,
        "Thank you",
        0.72,
        5.75,
        4.2,
        0.46,
        size=25,
        bold=True,
        color=BLUE,
        font=FONT_EN,
    )
    add_text(
        s,
        "FFW-SH5 · MuJoCo · ACT · PTE",
        0.74,
        6.29,
        5.6,
        0.28,
        size=12,
        color=MUTED,
        font=FONT_EN,
    )

    return new_slides


def validate(prs):
    assert len(prs.slides) == 15, f"Expected 15 slides, got {len(prs.slides)}"
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                raise ValueError(f"Negative shape location on slide {slide_idx}")
            if shape.left + shape.width > prs.slide_width + Inches(0.02):
                raise ValueError(
                    f"Shape exceeds slide width on slide {slide_idx}: {getattr(shape, 'name', '')}"
                )
            if shape.top + shape.height > prs.slide_height + Inches(0.02):
                raise ValueError(
                    f"Shape exceeds slide height on slide {slide_idx}: {getattr(shape, 'name', '')}"
                )


def main():
    rows = read_results()
    counts, samples, collector_episode = scan_dataset()
    required = {"green", "red", "orange", "blue"}
    if not required.issubset(samples):
        raise RuntimeError(f"Missing sample variants: {required - set(samples)}")
    build_assets(rows, counts, samples, collector_episode)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_cover(prs)
    append_main_slides(prs, rows, counts)
    validate(prs)
    prs.save(OUTPUT)

    reopened = Presentation(OUTPUT)
    validate(reopened)
    print(f"Created: {OUTPUT}")
    print(f"Slides: {len(reopened.slides)} (media-placeholder template)")
    print(f"Dataset distribution: {dict(counts)}")


if __name__ == "__main__":
    main()
